"""
genera_ppt.py — presentazione InfoVis
Intermedio (~3 min): codice D3 mostrato esplicitamente
Finale   (~7 min): solo schemi, flussi, tabelle — niente snippet
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SFONDO  = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT  = RGBColor(0x0F, 0x3A, 0x60)
AZZURRO = RGBColor(0x00, 0xB4, 0xD8)
BIANCO  = RGBColor(0xFF, 0xFF, 0xFF)
GRIGIO  = RGBColor(0xB0, 0xC4, 0xDE)
ROSSO   = RGBColor(0xE7, 0x4C, 0x3C)
VERDE   = RGBColor(0x27, 0xAE, 0x60)
GIALLO  = RGBColor(0xF5, 0xA6, 0x23)

W = Inches(13.33)
H = Inches(7.5)


# ─── Helper ──────────────────────────────────────────────────────────────────

def nuova_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def sfondo_slide(slide):
    shp = slide.shapes.add_shape(1, Inches(0), Inches(0), W, H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = SFONDO
    shp.line.fill.background()
    el = shp._element
    el.getparent().remove(el)
    slide.shapes._spTree.insert(2, el)

def rettangolo(slide, x, y, w, h, colore):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = colore
    shp.line.fill.background()
    return shp

def barra_titolo(slide, testo, size=26):
    shp = rettangolo(slide, Inches(0), Inches(0), W, Inches(1.05), ACCENT)
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = testo
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = AZZURRO

def sep_orizz(slide, y=Inches(1.15)):
    rettangolo(slide, Inches(0.6), y, Inches(12.13), Inches(0.04), AZZURRO)

def txt(slide, testo, x, y, w, h,
        size=15, bold=False, colore=None, align=PP_ALIGN.LEFT, italic=False):
    if colore is None:
        colore = GRIGIO
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = testo
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colore

def bullets(slide, voci, x, y, w, h, size=14, colore=None, marker="->"):
    if colore is None:
        colore = GRIGIO
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, voce in enumerate(voci):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = marker + "  " + voce
        run.font.size = Pt(size)
        run.font.color.rgb = colore

def codice(slide, testo, x, y, w, h, size=11):
    shp = rettangolo(slide, x, y, w, h, RGBColor(0x06, 0x0B, 0x1A))
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = testo
    run.font.size = Pt(size)
    run.font.color.rgb = AZZURRO

def tabella(slide, intestaz, righe, x, y, col_w, riga_h):
    x0, y0 = x, y
    for j, (h_txt, cw) in enumerate(zip(intestaz, col_w)):
        x_c = x0 + sum(col_w[:j])
        cella = rettangolo(slide, x_c, y0, cw, riga_h, ACCENT)
        tf = cella.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h_txt
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = AZZURRO
    for i, riga in enumerate(righe):
        bg_r = RGBColor(0x10, 0x10, 0x25) if i % 2 == 0 else SFONDO
        for j, (val, cw) in enumerate(zip(riga, col_w)):
            x_c = x0 + sum(col_w[:j])
            cella = rettangolo(slide, x_c, y0 + riga_h * (i + 1), cw, riga_h, bg_r)
            tf = cella.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(12)
            run.font.color.rgb = GRIGIO

def nodo(slide, testo, x, y, w, h, col=ACCENT, col_testo=BIANCO, size=12):
    """Box con barra laterale azzurra — usato per i flussi."""
    rettangolo(slide, x, y, w, h, col)
    rettangolo(slide, x, y, Inches(0.07), h, AZZURRO)
    txt(slide, testo, x + Inches(0.15), y + Inches(0.08),
        w - Inches(0.2), h - Inches(0.16), size=size, colore=col_testo)

def freccia_txt(slide, x, y):
    txt(slide, "->", x, y, Inches(0.45), Inches(0.4),
        size=17, bold=True, colore=AZZURRO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Copertina
# ════════════════════════════════════════════════════════════════════════════════

def slide_copertina(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    rettangolo(sl, Inches(0), Inches(0), Inches(0.35), H, AZZURRO)
    txt(sl, "Visualizzazione delle Informazioni — A.A. 2024/2025",
        Inches(0.7), Inches(1.8), Inches(12.2), Inches(0.6), size=18, colore=GRIGIO)
    txt(sl, "Progetto Intermedio e Progetto Finale",
        Inches(0.7), Inches(2.55), Inches(12.2), Inches(1.2),
        size=40, bold=True, colore=BIANCO)
    sep_orizz(sl, Inches(4.1))
    txt(sl, "D3.js v7  —  JavaScript ES6  —  HTML5 / CSS3",
        Inches(0.7), Inches(4.35), Inches(10), Inches(0.55), size=17, colore=AZZURRO)
    txt(sl, "Come D3 trasforma dati in visualizzazioni interattive.",
        Inches(0.7), Inches(5.05), Inches(12.0), Inches(0.65),
        size=14, italic=True, colore=GRIGIO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Separatore Intermedio
# ════════════════════════════════════════════════════════════════════════════════

def slide_sep_intermedio(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    rettangolo(sl, Inches(0), Inches(2.6), W, Inches(2.2), ACCENT)
    txt(sl, "PARTE 1  —  3 minuti", Inches(0), Inches(1.9), W, Inches(0.6),
        size=15, colore=AZZURRO, align=PP_ALIGN.CENTER)
    txt(sl, "Progetto Intermedio", Inches(0), Inches(2.65), W, Inches(1.1),
        size=44, bold=True, colore=BIANCO, align=PP_ALIGN.CENTER)
    txt(sl, "Glifi interattivi con D3.js: encode di 6 variabili su stick-figure",
        Inches(0), Inches(3.85), W, Inches(0.7), size=17, colore=GRIGIO,
        align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Intermedio: obiettivo + struttura
# ════════════════════════════════════════════════════════════════════════════════

def slide_intermedio_overview(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Intermedio — Obiettivo e Struttura")
    sep_orizz(sl)

    txt(sl, "Cosa fa", Inches(0.5), Inches(1.3), Inches(5.9), Inches(0.45),
        size=16, bold=True, colore=AZZURRO)
    bullets(sl, [
        "10 data-case rappresentati come stick-figure sul piano 2D",
        "Ogni glifo codifica 6 variabili: posizione X/Y + 4 via click",
        "Click cicla tra 3 coppie di variabili (var1/2, var3/4, var5/6)",
        "Transizione animata 800ms (d3.easeCubicInOut) al cambio stato",
        "Dominio globale: le 6 variabili condividono la stessa scala",
    ], Inches(0.5), Inches(1.8), Inches(5.9), Inches(3.0), size=14)

    txt(sl, "Struttura del progetto", Inches(7.1), Inches(1.3), Inches(5.8), Inches(0.45),
        size=16, bold=True, colore=AZZURRO)
    files = [
        ("index.html",        "struttura DOM, carica D3 e app.js"),
        ("data/dataset.json", "10 oggetti: id, name, var1 ... var6"),
        ("src/app.js",        "logica completa: scale, glifo, click"),
        ("src/styles.css",    "layout e tema visivo"),
    ]
    for i, (nome, desc) in enumerate(files):
        y_r = Inches(1.8) + i * Inches(0.75)
        bg_r = ACCENT if i % 2 == 0 else SFONDO
        rettangolo(sl, Inches(7.1), y_r, Inches(5.8), Inches(0.72), bg_r)
        txt(sl, nome, Inches(7.2), y_r + Inches(0.08), Inches(2.2), Inches(0.55),
            size=13, bold=True, colore=AZZURRO)
        txt(sl, desc, Inches(9.5), y_r + Inches(0.08), Inches(3.3), Inches(0.55),
            size=12, colore=GRIGIO)

    rettangolo(sl, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.8), ACCENT)
    txt(sl, "Serve HTTP server: d3.json() usa fetch() — bloccato da file:// per CORS policy.",
        Inches(0.65), Inches(5.18), Inches(12.0), Inches(0.65),
        size=13, italic=True, colore=GRIGIO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — I fondamentali D3 (con codice)
# ════════════════════════════════════════════════════════════════════════════════

def slide_intermedio_d3(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Intermedio — I Concetti D3 in Pratica")
    sep_orizz(sl)

    # sx: scale + data join
    txt(sl, "1. Scale lineari", Inches(0.5), Inches(1.3), Inches(5.9), Inches(0.4),
        size=14, bold=True, colore=AZZURRO)
    codice(sl,
        "var allVars = [];\n"
        "data.forEach(d => {\n"
        "  allVars.push(d.var1, d.var2, ..., d.var6);\n"
        "});\n"
        "var xScale = d3.scaleLinear()\n"
        "  .domain([d3.min(allVars), d3.max(allVars)])\n"
        "  .range([0, innerWidth]);\n"
        "var yScale = d3.scaleLinear()\n"
        "  .domain([d3.min(allVars), d3.max(allVars)])\n"
        "  .range([innerHeight, 0]);  // Y invertita",
        Inches(0.5), Inches(1.75), Inches(5.9), Inches(2.35))
    txt(sl, "Dominio globale: stesso min/max per X e Y -> confronto equo",
        Inches(0.5), Inches(4.18), Inches(5.9), Inches(0.4),
        size=12, italic=True, colore=GIALLO)

    txt(sl, "2. Data join — selectAll / data / enter / append",
        Inches(0.5), Inches(4.65), Inches(5.9), Inches(0.4),
        size=14, bold=True, colore=AZZURRO)
    codice(sl,
        "var omini = areaDisegno\n"
        "  .selectAll('.stickman-group')\n"
        "  .data(data)\n"
        "  .enter()\n"
        "  .append('g')\n"
        "  .attr('class', 'stickman-group')\n"
        "  .attr('transform', d =>\n"
        "    'translate('+xScale(d.var1)+','+yScale(d.var2)+')')",
        Inches(0.5), Inches(5.1), Inches(5.9), Inches(1.85))

    # dx: margin + transizione
    txt(sl, "3. Margin convention", Inches(7.1), Inches(1.3), Inches(5.8), Inches(0.4),
        size=14, bold=True, colore=AZZURRO)
    codice(sl,
        "var w = 900, h = 650;\n"
        "var margin = {top:50,right:50,\n"
        "              bottom:50,left:50};\n"
        "// area utile: 800 x 550 px\n\n"
        "var svg = d3.select('#chart')\n"
        "  .append('svg')\n"
        "  .attr('width', w).attr('height', h);\n\n"
        "var g = svg.append('g')\n"
        "  .attr('transform',\n"
        "    'translate('+margin.left+','+margin.top+')')",
        Inches(7.1), Inches(1.75), Inches(5.8), Inches(2.6))

    txt(sl, "4. Transizioni animate", Inches(7.1), Inches(4.45), Inches(5.8), Inches(0.4),
        size=14, bold=True, colore=AZZURRO)
    codice(sl,
        "omino.select('circle')\n"
        "  .transition()\n"
        "  .duration(800)\n"
        "  .ease(d3.easeCubicInOut)\n"
        "  .attr('transform',\n"
        "    'translate('+xScale(newX)+','+yScale(newY)+')')",
        Inches(7.1), Inches(4.9), Inches(5.8), Inches(1.75))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Glifo + logica click (con codice)
# ════════════════════════════════════════════════════════════════════════════════

def slide_intermedio_glifo(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Intermedio — Il Glifo e la Logica Click")
    sep_orizz(sl)

    txt(sl, "createStickman(gruppo, scale)  — anatomia SVG  (size = 25px con scale=1)",
        Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.45),
        size=15, bold=True, colore=AZZURRO)

    tabella(sl,
        ["Parte corpo", "Elemento SVG", "Coordinate"],
        [
            ["Testa",      "circle", "cy = -27.5px    r = 8.75px"],
            ["Occhi (x2)", "circle", "cx = +-2.5px    r = 1.25px    fill = #fff"],
            ["Corpo",      "line",   "y1 = -22.5px   ->   y2 = 5px"],
            ["Braccia",    "line",   "y = -7.5px    x = +-10px"],
            ["Gambe (x2)", "line",   "(0, 5px)   ->   (+-7.5px, 17.5px)"],
        ],
        Inches(0.5), Inches(1.82),
        [Inches(1.9), Inches(1.9), Inches(8.6)], Inches(0.5))

    txt(sl, "Logica click — 3 stati con operatore modulo",
        Inches(0.5), Inches(4.65), Inches(7.0), Inches(0.45),
        size=15, bold=True, colore=AZZURRO)
    codice(sl,
        "omini.on('click', function(event, d) {\n"
        "  d.clickState = (d.clickState + 1) % 3;\n"
        "  // 0 -> var1/var2,  1 -> var3/var4,  2 -> var5/var6\n"
        "  var newX = [d.var1, d.var3, d.var5][d.clickState];\n"
        "  var newY = [d.var2, d.var4, d.var6][d.clickState];\n"
        "  d3.select(this)\n"
        "    .transition().duration(800).ease(d3.easeCubicInOut)\n"
        "    .attr('transform',\n"
        "      'translate('+xScale(newX)+','+yScale(newY)+')');\n"
        "});",
        Inches(0.5), Inches(5.15), Inches(7.0), Inches(2.15))

    txt(sl, "Perche' clickState nel dato?",
        Inches(7.8), Inches(4.65), Inches(5.0), Inches(0.45),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "D3 lega i dati agli elementi SVG",
        "Ogni omino porta il suo stato con se'",
        "clickState = 0 aggiunto a runtime in createVisualization()",
        "Pattern D3: il dato e' la sorgente di verita' dello stato visivo",
        "Nessuna variabile globale necessaria",
    ], Inches(7.8), Inches(5.15), Inches(5.0), Inches(2.15), size=13, marker="*")


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Separatore Finale
# ════════════════════════════════════════════════════════════════════════════════

def slide_sep_finale(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    rettangolo(sl, Inches(0), Inches(2.6), W, Inches(2.2), ACCENT)
    txt(sl, "PARTE 2  —  7 minuti", Inches(0), Inches(1.9), W, Inches(0.6),
        size=15, colore=AZZURRO, align=PP_ALIGN.CENTER)
    txt(sl, "Progetto Finale", Inches(0), Inches(2.65), W, Inches(1.1),
        size=44, bold=True, colore=BIANCO, align=PP_ALIGN.CENTER)
    txt(sl, "Dashboard interattiva per Insider Threat Detection — CERT r4.2",
        Inches(0), Inches(3.85), W, Inches(0.7), size=17, colore=GRIGIO,
        align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Problema e dataset
# ════════════════════════════════════════════════════════════════════════════════

def slide_dataset(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — Il Problema e il Dataset")
    sep_orizz(sl)

    rettangolo(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.85), ACCENT)
    rettangolo(sl, Inches(0.5), Inches(1.3), Inches(0.1), Inches(0.85), ROSSO)
    txt(sl, "Un analista SOC deve identificare chi, tra 300 dipendenti, ha "
            "comportamenti anomali — senza scorrere 2400 righe CSV.",
        Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.65),
        size=15, bold=True, colore=BIANCO)

    txt(sl, "Struttura dati", Inches(0.5), Inches(2.35), Inches(5.9), Inches(0.4),
        size=15, bold=True, colore=AZZURRO)
    bullets(sl, [
        "300 utenti  x  8 settimane  =  2400 righe",
        "15 insider reali (5% del totale)",
        "Due CSV: profili utente + serie temporale settimanale",
        "Pipeline ML gia' eseguita: Isolation Forest + K-Means k=5",
        "La dashboard non fa ML: visualizza il risultato",
    ], Inches(0.5), Inches(2.82), Inches(5.9), Inches(3.2), size=14)

    txt(sl, "Le 8 feature principali", Inches(7.1), Inches(2.35), Inches(5.8), Inches(0.4),
        size=15, bold=True, colore=AZZURRO)
    feats = [
        ("final_anomaly_score", "score Isolation Forest"),
        ("rank",                "posizione in classifica rischio"),
        ("cluster",             "K-Means k=5, etichetta 0-4"),
        ("is_it_admin",         "flag 0/1"),
        ("work_hour_ratio",     "% accessi in orario"),
        ("after_hour_ratio",    "% accessi fuori orario"),
        ("file_ops",            "operazioni su file"),
        ("device_events",       "eventi su dispositivi USB"),
    ]
    for i, (feat, desc) in enumerate(feats):
        y_r = Inches(2.82) + i * Inches(0.52)
        bg_r = ACCENT if i % 2 == 0 else SFONDO
        rettangolo(sl, Inches(7.1), y_r, Inches(5.8), Inches(0.5), bg_r)
        txt(sl, feat, Inches(7.2), y_r + Inches(0.05), Inches(2.3), Inches(0.4),
            size=11, bold=True, colore=AZZURRO)
        txt(sl, desc, Inches(9.55), y_r + Inches(0.05), Inches(3.2), Inches(0.4),
            size=11, colore=GRIGIO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — I due file CSV: struttura e colonne
# ════════════════════════════════════════════════════════════════════════════════

def slide_dati_csv(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — I Dati: i due file CSV")
    sep_orizz(sl)

    # ── CSV 1 (sinistra) ──────────────────────────────────────────────────────
    rettangolo(sl, Inches(0.4), Inches(1.28), Inches(6.1), Inches(0.55),
               RGBColor(0x06, 0x0B, 0x1A))
    rettangolo(sl, Inches(0.4), Inches(1.28), Inches(0.1), Inches(0.55), AZZURRO)
    txt(sl, "anomalies_temporal_v2.csv  —  2400 righe  (300 utenti x 8 settimane)",
        Inches(0.6), Inches(1.35), Inches(5.75), Inches(0.42),
        size=12, bold=True, colore=AZZURRO)

    cols_temporal = [
        ("user_id",              "stringa",  "identificatore utente"),
        ("week",                 "int 1-8",  "settimana di osservazione"),
        ("cluster",              "int 0-4",  "cluster K-Means assegnato"),
        ("insider",              "0 / 1",    "etichetta reale: 1 = insider confermato"),
        ("rank",                 "int",      "posizione classifica rischio (1 = piu' anomalo)"),
        ("final_anomaly_score",  "float",    "score Isolation Forest (puo' essere negativo)"),
        ("n_logon / n_usb / n_file / n_email / n_http", "int", "conteggi attivita' per tipo"),
        ("n_afterhourallact",    "int",      "attivita' totali fuori orario lavorativo"),
        ("n_allact",             "int",      "attivita' totali"),
        ("O, C, E, A, N",        "float",    "Big Five personality traits"),
        ("cluster_distance",     "float",    "distanza dal centroide del proprio cluster"),
    ]
    for i, (col, tipo, desc) in enumerate(cols_temporal):
        y_r = Inches(1.9) + i * Inches(0.44)
        bg_r = ACCENT if i % 2 == 0 else RGBColor(0x10, 0x10, 0x25)
        rettangolo(sl, Inches(0.4), y_r, Inches(6.1), Inches(0.42), bg_r)
        txt(sl, col,  Inches(0.5),  y_r + Inches(0.05), Inches(2.3), Inches(0.33), size=10, bold=True, colore=AZZURRO)
        txt(sl, tipo, Inches(2.85), y_r + Inches(0.05), Inches(0.9), Inches(0.33), size=10, colore=GIALLO)
        txt(sl, desc, Inches(3.8),  y_r + Inches(0.05), Inches(2.6), Inches(0.33), size=10, colore=GRIGIO)

    # ── CSV 2 (destra) ───────────────────────────────────────────────────────
    rettangolo(sl, Inches(6.85), Inches(1.28), Inches(6.1), Inches(0.55),
               RGBColor(0x06, 0x0B, 0x1A))
    rettangolo(sl, Inches(6.85), Inches(1.28), Inches(0.1), Inches(0.55), VERDE)
    txt(sl, "cluster_profiles_v2.csv  —  5 righe  (una per cluster)",
        Inches(7.05), Inches(1.35), Inches(5.75), Inches(0.42),
        size=12, bold=True, colore=VERDE)

    cols_profiles = [
        ("cluster",           "int 0-4",  "identificatore cluster"),
        ("n_logon",           "float",    "media logon nel cluster"),
        ("n_usb",             "float",    "media eventi USB"),
        ("n_file",            "float",    "media operazioni su file"),
        ("n_email",           "float",    "media email inviate"),
        ("n_http",            "float",    "media richieste HTTP"),
        ("n_afterhourallact", "float",    "media attivita' fuori orario"),
        ("final_anomaly_score","float",   "media score Isolation Forest"),
        ("cluster_distance",  "float",    "distanza media dal centroide"),
        ("O, C, E, A, N",     "float",    "media Big Five per cluster"),
    ]
    for i, (col, tipo, desc) in enumerate(cols_profiles):
        y_r = Inches(1.9) + i * Inches(0.44)
        bg_r = ACCENT if i % 2 == 0 else RGBColor(0x10, 0x10, 0x25)
        rettangolo(sl, Inches(6.85), y_r, Inches(6.1), Inches(0.42), bg_r)
        txt(sl, col,  Inches(6.95), y_r + Inches(0.05), Inches(2.3), Inches(0.33), size=10, bold=True, colore=VERDE)
        txt(sl, tipo, Inches(9.3),  y_r + Inches(0.05), Inches(0.9), Inches(0.33), size=10, colore=GIALLO)
        txt(sl, desc, Inches(10.25),y_r + Inches(0.05), Inches(2.6), Inches(0.33), size=10, colore=GRIGIO)

    # ── nota chi usa cosa ────────────────────────────────────────────────────
    rettangolo(sl, Inches(0.4), Inches(6.7), Inches(12.55), Inches(0.6), ACCENT)
    rettangolo(sl, Inches(0.4), Inches(6.7), Inches(0.08), Inches(0.6), AZZURRO)
    txt(sl, "anomalies_temporal_v2  ->  univariate, bivariate, trivariate, temporal, parallel-coord  (dati filtrabili)   |   "
            "cluster_profiles_v2  ->  radar + heatmap  (medie fisse per cluster)",
        Inches(0.6), Inches(6.77), Inches(12.2), Inches(0.48),
        size=11, bold=True, colore=GRIGIO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Architettura: schema flusso moduli
# ════════════════════════════════════════════════════════════════════════════════

def slide_architettura(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — Architettura del Codice")
    sep_orizz(sl)

    # layer INPUT
    txt(sl, "INPUT", Inches(0.5), Inches(1.35), Inches(2.5), Inches(0.35),
        size=11, bold=True, colore=AZZURRO)
    nodo(sl, "data/\n*.csv\n(2 file)",
         Inches(0.5), Inches(1.75), Inches(2.3), Inches(1.2),
         col=RGBColor(0x06, 0x0B, 0x1A), col_testo=AZZURRO, size=12)
    freccia_txt(sl, Inches(2.9), Inches(2.15))

    # layer CORE
    txt(sl, "CORE", Inches(3.35), Inches(1.35), Inches(3.0), Inches(0.35),
        size=11, bold=True, colore=AZZURRO)
    core_mod = [
        ("config.js",       "colori, label, soglie"),
        ("data-loader.js",  "carica + filtra dati"),
        ("interactions.js", "ascolta filtri UI"),
        ("main.js",         "bootstrap + coordinamento"),
    ]
    for i, (nome, desc) in enumerate(core_mod):
        y_n = Inches(1.75) + i * Inches(0.72)
        nodo(sl, nome + "\n" + desc,
             Inches(3.35), y_n, Inches(2.9), Inches(0.65), size=11)
    freccia_txt(sl, Inches(6.35), Inches(2.15))

    # layer VISUALIZZAZIONI
    txt(sl, "VISUALIZZAZIONI  (init + update)", Inches(6.8), Inches(1.35),
        Inches(6.0), Inches(0.35), size=11, bold=True, colore=AZZURRO)
    viz_mod = [
        ("univariate.js",           "#istogramma"),
        ("bivariate.js",            "#scatter-rank-score"),
        ("trivariate.js",           "#scatter-multivariato"),
        ("temporal.js",             "#grafico-temporale"),
        ("parallel-coordinates.js", "#coordinate-parallele"),
        ("multivariate.js",         "#radar  +  #heatmap"),
    ]
    for i, (nome, dom) in enumerate(viz_mod):
        y_n = Inches(1.75) + i * Inches(0.72)
        nodo(sl, nome, Inches(6.8), y_n, Inches(2.5), Inches(0.65), size=11)
        txt(sl, dom, Inches(9.4), y_n + Inches(0.15),
            Inches(3.2), Inches(0.4), size=11, italic=True, colore=GIALLO)

    # nota interfaccia comune
    rettangolo(sl, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.65), ACCENT)
    txt(sl, "Interfaccia comune a tutti i moduli visualizzazione:  "
            "init(svgId, data)   |   update(filteredData)",
        Inches(0.65), Inches(6.62), Inches(12.0), Inches(0.5),
        size=13, bold=True, colore=AZZURRO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Architettura dati → grafici (schema colonne)
# ════════════════════════════════════════════════════════════════════════════════

def slide_architettura_dati(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — Quale Colonna Entra in Quale Grafico")
    sep_orizz(sl)

    txt(sl, "anomalies_temporal_v2.csv  (righe per utente x settimana)",
        Inches(0.4), Inches(1.28), Inches(8.0), Inches(0.4),
        size=13, bold=True, colore=AZZURRO)
    txt(sl, "cluster_profiles_v2.csv  (medie per cluster)",
        Inches(8.6), Inches(1.28), Inches(4.4), Inches(0.4),
        size=13, bold=True, colore=VERDE)

    # tabella: colonna | grafico che la usa | encoding
    righe_mappa = [
        ("final_anomaly_score", "Istogramma, Scatter, Bubble, Line, Parallel", "asse X / Y / colore"),
        ("rank",                "Scatter rank-score, Parallel",                "asse X / asse"),
        ("insider",             "Scatter rank-score, Bubble",                  "colore binario rosso/verde"),
        ("cluster",             "Bubble, Line, Parallel",                      "colore ordinale (schemeCategory10)"),
        ("n_afterhourallact",   "Bubble (asse X), Parallel (asse)",            "posizione X / asse verticale"),
        ("n_allact",            "Bubble (dimensione cerchio)",                 "scaleSqrt -> raggio"),
        ("week",                "Line chart",                                  "asse X temporale (1-8)"),
        ("n_logon/usb/file/email/http", "Parallel, Radar*, Heatmap*",         "assi paralleli / raggi / celle"),
        ("O, C, E, A, N",       "Parallel, Heatmap*",                         "assi paralleli / celle"),
        ("cluster_distance",    "Heatmap*",                                    "cella feature"),
    ]
    col_w_m = [Inches(3.2), Inches(5.4), Inches(3.85)]
    x0, y0, rh = Inches(0.4), Inches(1.75), Inches(0.46)

    for j, (h, cw) in enumerate(zip(["Colonna CSV", "Grafico(i) che la usano", "Encoding D3"], col_w_m)):
        x_c = x0 + sum(col_w_m[:j])
        cella = rettangolo(sl, x_c, y0, cw, rh, ACCENT)
        tf = cella.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = AZZURRO

    for i, (col, grafici, enc) in enumerate(righe_mappa):
        bg_r = RGBColor(0x10, 0x10, 0x25) if i % 2 == 0 else SFONDO
        col_testo = VERDE if "*" in grafici else GRIGIO
        grafici_clean = grafici.replace("*", "")
        for j, (val, cw) in enumerate(zip([col, grafici_clean, enc], col_w_m)):
            x_c = x0 + sum(col_w_m[:j])
            cella = rettangolo(sl, x_c, y0 + rh * (i + 1), cw, rh, bg_r)
            tf = cella.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = val
            run.font.size = Pt(10)
            run.font.color.rgb = (AZZURRO if j == 0 else col_testo if j == 1 else GIALLO)

    rettangolo(sl, Inches(0.4), Inches(6.6), Inches(12.55), Inches(0.6), ACCENT)
    txt(sl, "* Radar e Heatmap leggono cluster_profiles_v2.csv direttamente — non usano filteredData",
        Inches(0.6), Inches(6.67), Inches(12.2), Inches(0.48),
        size=11, italic=True, colore=VERDE)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Grafico 1: Istogramma
# ════════════════════════════════════════════════════════════════════════════════

def slide_graf_istogramma(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Grafico 1 — Istogramma  (univariate.js)")
    sep_orizz(sl)

    # header badges
    for i, (label, col) in enumerate([("Univariata", AZZURRO), ("filteredData", VERDE), ("scaleLinear", GIALLO)]):
        x_b = Inches(0.5) + i * Inches(2.55)
        rettangolo(sl, x_b, Inches(1.28), Inches(2.3), Inches(0.42), ACCENT)
        txt(sl, label, x_b + Inches(0.1), Inches(1.31), Inches(2.1), Inches(0.36),
            size=12, bold=True, colore=col)

    # domanda
    rettangolo(sl, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.55), RGBColor(0x06,0x0B,0x1A))
    txt(sl, "Domanda: come sono distribuiti gli anomaly score? Ci sono due popolazioni distinte?",
        Inches(0.65), Inches(1.89), Inches(12.0), Inches(0.42),
        size=13, italic=True, colore=BIANCO)

    # sx: dati → scala → SVG
    txt(sl, "Flusso dati -> SVG", Inches(0.5), Inches(2.55), Inches(6.0), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    flusso = [
        ("INPUT",   "filteredData  (array di oggetti utente)"),
        ("ESTRAE",  "data.map(d => d.final_anomaly_score)  — solo questa colonna"),
        ("BIN",     "d3.bin()  con 20 bin sull'extent dei valori"),
        ("SCALA X", "scaleLinear  domain=[min,max]  range=[0, width]"),
        ("SCALA Y", "scaleLinear  domain=[0, maxCount]  range=[height, 0]  (Y invertita)"),
        ("SVG",     "un <rect> per bin  — altezza proporzionale al conteggio"),
    ]
    for i, (label, desc) in enumerate(flusso):
        y_f = Inches(3.0) + i * Inches(0.6)
        rettangolo(sl, Inches(0.5), y_f, Inches(6.0), Inches(0.55), ACCENT if i%2==0 else SFONDO)
        rettangolo(sl, Inches(0.5), y_f, Inches(0.06), Inches(0.55), AZZURRO)
        txt(sl, label, Inches(0.7), y_f+Inches(0.07), Inches(1.1), Inches(0.4), size=11, bold=True, colore=AZZURRO)
        txt(sl, desc,  Inches(1.85), y_f+Inches(0.07), Inches(4.55), Inches(0.4), size=11, colore=GRIGIO)

    # dx: interazione + scelte tecniche
    txt(sl, "Interazione e scelte tecniche", Inches(7.0), Inches(2.55), Inches(5.8), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "mouseover: barra diventa rossa (accent) + tooltip con range bin e conteggio",
        "mouseout: ripristina colore primary (blu)",
        "Animazione: barre partono da altezza 0 e crescono in 750ms",
        "update(): rimuove il vecchio SVG e ridisegna da zero",
        "Scelta bin = 20: abbastanza dettaglio senza rumore",
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.2), size=13)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Grafico 2: Scatter rank vs score
# ════════════════════════════════════════════════════════════════════════════════

def slide_graf_scatter(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Grafico 2 — Scatter rank vs score  (bivariate.js)")
    sep_orizz(sl)

    for i, (label, col) in enumerate([("Bivariata", AZZURRO), ("filteredData", VERDE), ("scaleLinear x2", GIALLO)]):
        x_b = Inches(0.5) + i * Inches(2.55)
        rettangolo(sl, x_b, Inches(1.28), Inches(2.3), Inches(0.42), ACCENT)
        txt(sl, label, x_b+Inches(0.1), Inches(1.31), Inches(2.1), Inches(0.36),
            size=12, bold=True, colore=col)

    rettangolo(sl, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.55), RGBColor(0x06,0x0B,0x1A))
    txt(sl, "Domanda: c'e' relazione tra posizione in classifica e score? Gli insider sono tutti in cima?",
        Inches(0.65), Inches(1.89), Inches(12.0), Inches(0.42),
        size=13, italic=True, colore=BIANCO)

    txt(sl, "Aggregazione per utente (fix dataset temporale)", Inches(0.5), Inches(2.55), Inches(6.0), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    flusso_b = [
        ("PROBLEMA",  "CSV ha 2400 righe: 300 utenti x 8 settimane -> rank diverso ogni settimana"),
        ("FIX",       "d3.rollup(data, media, user_id) -> 300 punti unici (un punto per utente)"),
        ("RANK",      "media dei rank nelle 8 settimane"),
        ("SCORE",     "media del final_anomaly_score nelle 8 settimane"),
        ("INSIDER",   "valore fisso (non cambia nel tempo)"),
    ]
    for i, (label, desc) in enumerate(flusso_b):
        y_f = Inches(3.0) + i * Inches(0.58)
        rettangolo(sl, Inches(0.5), y_f, Inches(6.0), Inches(0.52), ACCENT if i%2==0 else SFONDO)
        rettangolo(sl, Inches(0.5), y_f, Inches(0.06), Inches(0.52), AZZURRO if i>0 else ROSSO)
        txt(sl, label, Inches(0.7), y_f+Inches(0.07), Inches(1.2), Inches(0.38), size=11, bold=True, colore=AZZURRO if i>0 else ROSSO)
        txt(sl, desc,  Inches(2.0), y_f+Inches(0.07), Inches(4.4), Inches(0.38), size=11, colore=GRIGIO)

    txt(sl, "Interazione e scelte tecniche", Inches(7.0), Inches(2.55), Inches(5.8), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "Ogni punto = un utente unico (300 punti, media 8 settimane)",
        "Encoding: X = rank medio, Y = score medio, Colore = insider (rosso/verde)",
        "mouseover: raggio 4px -> 8px + tooltip (user_id, rank medio, score medio, tipo)",
        "Animazione entrata: raggio 0 -> 4px con delay i*5ms (effetto cascata)",
        "Legenda manuale in alto a destra: cerchio rosso Insider / verde Normale",
        "Colore binario: piu' leggibile di scala continua per variabile 0/1",
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.2), size=13)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Grafico 3: Bubble chart
# ════════════════════════════════════════════════════════════════════════════════

def slide_graf_bubble(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Grafico 3 — Bubble chart  (trivariate.js)")
    sep_orizz(sl)

    for i, (label, col) in enumerate([("Trivariata", AZZURRO), ("filteredData", VERDE), ("scaleSqrt", GIALLO)]):
        x_b = Inches(0.5) + i * Inches(2.55)
        rettangolo(sl, x_b, Inches(1.28), Inches(2.3), Inches(0.42), ACCENT)
        txt(sl, label, x_b+Inches(0.1), Inches(1.31), Inches(2.1), Inches(0.36),
            size=12, bold=True, colore=col)

    rettangolo(sl, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.55), RGBColor(0x06,0x0B,0x1A))
    txt(sl, "Domanda: chi lavora troppo fuori orario? Quanto e' attivo in totale? Di che cluster fa parte?",
        Inches(0.65), Inches(1.89), Inches(12.0), Inches(0.42),
        size=13, italic=True, colore=BIANCO)

    txt(sl, "Aggregazione per utente (fix dataset temporale)", Inches(0.5), Inches(2.55), Inches(6.0), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    flusso_bub = [
        ("PROBLEMA",  "CSV 2400 righe: 300 utenti x 8 settimane -> dati diversi per settimana"),
        ("FIX",       "d3.rollup(data, media, user_id) -> 300 bolle uniche per utente"),
        ("X",         "media di n_afterhourallact sulle 8 settimane"),
        ("Y / SIZE",  "media di final_anomaly_score / n_allact sulle 8 settimane"),
        ("CLUSTER",   "valore fisso per utente (non cambia nel tempo)"),
    ]
    for i, (label, desc) in enumerate(flusso_bub):
        y_f = Inches(3.0) + i * Inches(0.55)
        rettangolo(sl, Inches(0.5), y_f, Inches(6.0), Inches(0.49), ACCENT if i%2==0 else SFONDO)
        rettangolo(sl, Inches(0.5), y_f, Inches(0.06), Inches(0.49), AZZURRO if i>0 else ROSSO)
        txt(sl, label, Inches(0.7), y_f+Inches(0.07), Inches(1.2), Inches(0.36), size=11, bold=True, colore=AZZURRO if i>0 else ROSSO)
        txt(sl, desc,  Inches(2.0), y_f+Inches(0.07), Inches(4.4), Inches(0.36), size=11, colore=GRIGIO)

    txt(sl, "Encoding e scelte tecniche", Inches(7.0), Inches(2.55), Inches(5.8), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "4 encoding: X = after-hour, Y = score, Colore = cluster, Size = tot. attivita'",
        "scaleSqrt per la dimensione: area visiva proporzionale al valore (Cleveland & McGill)",
        "Due legende a destra: cluster (cerchi colorati) + dimensione (3 taglie)",
        "Larghezza SVG da getBoundingClientRect() -> si adatta alla griglia CSS",
        "Tooltip: user_id, after-hour medio, score medio, cluster, tot. attivita' media, tipo",
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.2), size=13)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Grafico 4: Line chart temporale
# ════════════════════════════════════════════════════════════════════════════════

def slide_graf_line(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Grafico 4 — Line chart temporale  (temporal.js)")
    sep_orizz(sl)

    for i, (label, col) in enumerate([("Temporale", AZZURRO), ("filteredData", VERDE), ("d3.rollup", GIALLO)]):
        x_b = Inches(0.5) + i * Inches(2.55)
        rettangolo(sl, x_b, Inches(1.28), Inches(2.3), Inches(0.42), ACCENT)
        txt(sl, label, x_b+Inches(0.1), Inches(1.31), Inches(2.1), Inches(0.36),
            size=12, bold=True, colore=col)

    rettangolo(sl, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.55), RGBColor(0x06,0x0B,0x1A))
    txt(sl, "Domanda: il rischio di un cluster aumenta nel tempo? C'e' una settimana critica?",
        Inches(0.65), Inches(1.89), Inches(12.0), Inches(0.42),
        size=13, italic=True, colore=BIANCO)

    txt(sl, "Aggregazione: da 2400 righe a 5 linee", Inches(0.5), Inches(2.55), Inches(6.0), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    agg = [
        ("INPUT",    "2400 righe  (300 utenti x 8 settimane)"),
        ("RAGGRUPPA","d3.rollup(data, mean(score), cluster, week)"),
        ("RISULTATO","Map annidata: cluster -> settimana -> score medio"),
        ("LINEE",    "5 path SVG  (una per cluster), 8 punti ciascuna"),
        ("ASSE X",   "scaleLinear  domain=[1,8]  — le settimane"),
        ("ASSE Y",   "scaleLinear  domain=[0, max_score]"),
    ]
    for i, (label, desc) in enumerate(agg):
        y_f = Inches(3.0) + i * Inches(0.57)
        rettangolo(sl, Inches(0.5), y_f, Inches(6.0), Inches(0.52), ACCENT if i%2==0 else SFONDO)
        rettangolo(sl, Inches(0.5), y_f, Inches(0.06), Inches(0.52), AZZURRO)
        txt(sl, label, Inches(0.7), y_f+Inches(0.07), Inches(1.05), Inches(0.38), size=11, bold=True, colore=AZZURRO)
        txt(sl, desc,  Inches(1.8), y_f+Inches(0.07), Inches(4.6), Inches(0.38), size=11, colore=GRIGIO)

    txt(sl, "Animazione stroke-dasharray", Inches(7.0), Inches(2.55), Inches(5.8), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "Misura la lunghezza totale della path SVG (getTotalLength)",
        "Imposta stroke-dasharray = L  e  stroke-dashoffset = L  (linea invisibile)",
        "Transizione: dashoffset -> 0  in 1500ms  ease(d3.easeLinear)",
        "Effetto: la linea si 'disegna' da sinistra a destra",
        "Punti circolari appaiono dopo la linea (delay 1500ms)",
        "C4 in rosso e grassetto nella legenda",
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.5), size=13)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Grafico 5: Parallel coordinates
# ════════════════════════════════════════════════════════════════════════════════

def slide_graf_parallel(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Grafico 5 — Parallel coordinates  (parallel-coordinates.js)")
    sep_orizz(sl)

    for i, (label, col) in enumerate([("Multivariata 6D", AZZURRO), ("filteredData", VERDE), ("scalePoint", GIALLO)]):
        x_b = Inches(0.5) + i * Inches(2.55)
        rettangolo(sl, x_b, Inches(1.28), Inches(2.3), Inches(0.42), ACCENT)
        txt(sl, label, x_b+Inches(0.1), Inches(1.31), Inches(2.1), Inches(0.36),
            size=12, bold=True, colore=col)

    rettangolo(sl, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.55), RGBColor(0x06,0x0B,0x1A))
    txt(sl, "Domanda: quale cluster e' outlier su piu' dimensioni contemporaneamente?",
        Inches(0.65), Inches(1.89), Inches(12.0), Inches(0.42),
        size=13, italic=True, colore=BIANCO)

    txt(sl, "I 6 assi e le scale usate", Inches(0.5), Inches(2.55), Inches(6.0), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    assi = [
        ("final_anomaly_score", "scaleLinear", "asse 1 — score Isolation Forest"),
        ("n_allact",            "scaleLinear", "asse 2 — attivita' totali"),
        ("n_afterhourallact",   "scaleLinear", "asse 3 — attivita' fuori orario"),
        ("n_email",             "scaleLinear", "asse 4 — email inviate"),
        ("n_file",              "scaleLinear", "asse 5 — operazioni su file"),
        ("rank",                "scaleLinear", "asse 6 — classifica rischio"),
    ]
    for i, (col_csv, scala, desc) in enumerate(assi):
        y_a = Inches(3.0) + i * Inches(0.57)
        rettangolo(sl, Inches(0.5), y_a, Inches(6.0), Inches(0.52), ACCENT if i%2==0 else SFONDO)
        rettangolo(sl, Inches(0.5), y_a, Inches(0.06), Inches(0.52), AZZURRO)
        txt(sl, col_csv, Inches(0.7), y_a+Inches(0.07), Inches(2.2), Inches(0.38), size=11, bold=True, colore=AZZURRO)
        txt(sl, scala,   Inches(2.95), y_a+Inches(0.07), Inches(1.1), Inches(0.38), size=11, colore=GIALLO)
        txt(sl, desc,    Inches(4.1),  y_a+Inches(0.07), Inches(2.3), Inches(0.38), size=11, colore=GRIGIO)

    txt(sl, "Scelte tecniche", Inches(7.0), Inches(2.55), Inches(5.8), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "scalePoint per l'asse X: posiziona i 6 assi equidistanti sulla larghezza",
        "Una scaleLinear separata per ogni asse (range valori molto diversi tra loro)",
        "Campionamento a 200 utenti se dataset > 200 (shuffle casuale .slice(0,200))",
        "2 layer: linee grigie opache (contesto) + linee cluster colorate (focus)",
        "mouseover: linea si ispessisce + .raise() la porta in primo piano",
        "Nota testuale in basso indica quanti utenti sono visualizzati",
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.5), size=13)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Grafico 6: Radar chart
# ════════════════════════════════════════════════════════════════════════════════

def slide_graf_radar(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Grafico 6 — Radar chart  (multivariate.js)")
    sep_orizz(sl)

    for i, (label, col) in enumerate([("Multivariata 6D", AZZURRO), ("cluster_profiles_v2.csv", VERDE), ("d3.lineRadial", GIALLO)]):
        x_b = Inches(0.5) + i * Inches(2.85)
        rettangolo(sl, x_b, Inches(1.28), Inches(2.6), Inches(0.42), ACCENT)
        txt(sl, label, x_b+Inches(0.1), Inches(1.31), Inches(2.4), Inches(0.36),
            size=12, bold=True, colore=col)

    rettangolo(sl, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.55), RGBColor(0x06,0x0B,0x1A))
    txt(sl, "Domanda: qual e' la 'forma' comportamentale di ogni cluster? C4 assomiglia agli altri?",
        Inches(0.65), Inches(1.89), Inches(12.0), Inches(0.42),
        size=13, italic=True, colore=BIANCO)

    txt(sl, "I 6 assi del radar (da cluster_profiles_v2.csv)", Inches(0.5), Inches(2.55), Inches(6.0), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    assi_r = [
        ("n_logon",           "media logon nel cluster"),
        ("n_usb",             "media eventi USB"),
        ("n_file",            "media operazioni su file"),
        ("n_email",           "media email inviate"),
        ("n_http",            "media richieste HTTP"),
        ("n_afterhourallact", "media attivita' fuori orario"),
    ]
    for i, (col_csv, desc) in enumerate(assi_r):
        y_a = Inches(3.0) + i * Inches(0.57)
        rettangolo(sl, Inches(0.5), y_a, Inches(6.0), Inches(0.52), ACCENT if i%2==0 else SFONDO)
        rettangolo(sl, Inches(0.5), y_a, Inches(0.06), Inches(0.52), AZZURRO)
        txt(sl, col_csv, Inches(0.7), y_a+Inches(0.07), Inches(2.3), Inches(0.38), size=11, bold=True, colore=AZZURRO)
        txt(sl, desc,    Inches(3.1), y_a+Inches(0.07), Inches(3.3), Inches(0.38), size=11, colore=GRIGIO)

    txt(sl, "Scelte tecniche chiave", Inches(7.0), Inches(2.55), Inches(5.8), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "Fonte: cluster_profiles_v2.csv — 5 righe (medie gia' calcolate, non filtrabili)",
        "Normalizzazione per-feature tra cluster: extent sui soli 5 valori -> [0.1, 0.9]",
        "Senza normalizzazione tutti i poligoni sarebbero sovrapposti vicino al centro",
        "d3.lineRadial() con curveLinearClosed: chiude automaticamente il poligono",
        "Struttura statica (griglia, assi) creata solo al primo render",
        "Poligoni colorati con stessa palette cluster degli altri grafici",
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.5), size=13)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Grafico 7: Heatmap
# ════════════════════════════════════════════════════════════════════════════════

def slide_graf_heatmap(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Grafico 7 — Heatmap  (multivariate.js)")
    sep_orizz(sl)

    for i, (label, col) in enumerate([("Multivariata 11D", AZZURRO), ("cluster_profiles_v2.csv", VERDE), ("scaleBand", GIALLO)]):
        x_b = Inches(0.5) + i * Inches(2.85)
        rettangolo(sl, x_b, Inches(1.28), Inches(2.6), Inches(0.42), ACCENT)
        txt(sl, label, x_b+Inches(0.1), Inches(1.31), Inches(2.4), Inches(0.36),
            size=12, bold=True, colore=col)

    rettangolo(sl, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.55), RGBColor(0x06,0x0B,0x1A))
    txt(sl, "Domanda: quale feature distingue di piu' i cluster? Dove C4 e' davvero diverso?",
        Inches(0.65), Inches(1.89), Inches(12.0), Inches(0.42),
        size=13, italic=True, colore=BIANCO)

    txt(sl, "Le 11 colonne nelle celle (da cluster_profiles_v2.csv)", Inches(0.5), Inches(2.55), Inches(6.0), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    cols_h = [
        ("O, C, E, A, N",        "Big Five personality traits (5 colonne)"),
        ("n_logon",              "media logon"),
        ("n_email",              "media email"),
        ("n_http",               "media HTTP"),
        ("n_afterhourallact",    "media attivita' fuori orario"),
        ("final_anomaly_score",  "media score Isolation Forest"),
        ("cluster_distance",     "distanza media dal centroide"),
    ]
    for i, (col_csv, desc) in enumerate(cols_h):
        y_a = Inches(3.0) + i * Inches(0.52)
        rettangolo(sl, Inches(0.5), y_a, Inches(6.0), Inches(0.47), ACCENT if i%2==0 else SFONDO)
        rettangolo(sl, Inches(0.5), y_a, Inches(0.06), Inches(0.47), AZZURRO)
        txt(sl, col_csv, Inches(0.7), y_a+Inches(0.06), Inches(2.3), Inches(0.35), size=10, bold=True, colore=AZZURRO)
        txt(sl, desc,    Inches(3.1), y_a+Inches(0.06), Inches(3.3), Inches(0.35), size=10, colore=GRIGIO)

    txt(sl, "Scelte tecniche chiave", Inches(7.0), Inches(2.55), Inches(5.8), Inches(0.38),
        size=14, bold=True, colore=AZZURRO)
    bullets(sl, [
        "scaleBand per righe (cluster) e colonne (feature): celle di larghezza fissa",
        "Normalizzazione per-colonna: min/max su tutti i cluster -> [0,1]",
        "Scala colore: scaleSequential(interpolateRdYlGn) — rosso basso, verde alto",
        "C4 ha bordo rosso (stroke: #e74c3c) per evidenziarlo visivamente",
        "mouseover: tooltip con valore reale + posizione percentuale relativa",
        "Legenda colore: gradiente pixel per pixel (loop su 100 step)",
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.5), size=13)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Flusso interno di ogni visualizzazione
# ════════════════════════════════════════════════════════════════════════════════

def slide_flusso_viz(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — Come Funziona Ogni Visualizzazione")
    sep_orizz(sl)

    txt(sl, "Ogni modulo segue lo stesso schema interno:",
        Inches(0.5), Inches(1.28), Inches(12.3), Inches(0.4),
        size=13, italic=True, colore=GRIGIO)

    steps = [
        "Dati grezzi\n(array JS)",
        "Scale D3\ndomain -> range",
        "Elementi SVG\nrect / circle\nline / path",
        "Transizione\nanimata",
        "DOM\naggiornato",
    ]
    colori_s = [ACCENT, ACCENT, ACCENT, ACCENT, RGBColor(0x0A, 0x3D, 0x2B)]
    for i, (label, col) in enumerate(zip(steps, colori_s)):
        x_n = Inches(0.4) + i * Inches(2.5)
        nodo(sl, label, x_n, Inches(1.82), Inches(2.2), Inches(1.0),
             col=col, size=12)
        if i < 4:
            freccia_txt(sl, x_n + Inches(2.25), Inches(2.1))

    # 4 box scale
    txt(sl, "Le scale usate nei 7 grafici:", Inches(0.5), Inches(3.1),
        Inches(12.3), Inches(0.4), size=14, bold=True, colore=AZZURRO)

    scale_dati = [
        ("scaleLinear",
         "Istogramma, Scatter\nX e Y su valori continui — relazione lineare dato/pixel"),
        ("scaleSqrt",
         "Bubble chart\nArea visiva proporzionale al valore, non al raggio (Cleveland & McGill)"),
        ("scalePoint",
         "Parallel coordinates\nDistribuisce assi discreti equidistanti sulla larghezza"),
        ("scaleBand",
         "Heatmap\nCelle di larghezza fissa per cluster x feature"),
    ]
    for i, (scala, desc) in enumerate(scale_dati):
        col_x = Inches(0.5) + (i % 2) * Inches(6.45)
        row_y = Inches(3.6) + (i // 2) * Inches(1.6)
        rettangolo(sl, col_x, row_y, Inches(6.0), Inches(1.45), ACCENT)
        rettangolo(sl, col_x, row_y, Inches(0.07), Inches(1.45), AZZURRO)
        txt(sl, scala, col_x + Inches(0.18), row_y + Inches(0.1),
            Inches(5.7), Inches(0.45), size=14, bold=True, colore=AZZURRO)
        txt(sl, desc, col_x + Inches(0.18), row_y + Inches(0.62),
            Inches(5.7), Inches(0.72), size=12, colore=GRIGIO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Flusso aggregazione dati + animazione linee
# ════════════════════════════════════════════════════════════════════════════════

def slide_flusso_temporale(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — Flusso: Aggregazione e Animazione")
    sep_orizz(sl)

    txt(sl, "Line chart — da 2400 righe a 5 linee animate:",
        Inches(0.5), Inches(1.28), Inches(12.3), Inches(0.4),
        size=14, bold=True, colore=AZZURRO)

    steps1 = [
        "2400 righe CSV\n(utente x settimana)",
        "Aggrega per\ncluster + settimana\n(media score)",
        "5 serie\n(una per cluster,\n8 punti ciascuna)",
        "5 path SVG\ncurveMonotoneX",
        "Animazione\nstroke-dasharray",
    ]
    for i, label in enumerate(steps1):
        col = RGBColor(0x0A, 0x3D, 0x2B) if i == 4 else ACCENT
        x_n = Inches(0.4) + i * Inches(2.5)
        nodo(sl, label, x_n, Inches(1.82), Inches(2.2), Inches(1.05), col=col, size=11)
        if i < 4:
            freccia_txt(sl, x_n + Inches(2.25), Inches(2.1))

    txt(sl, "Come funziona stroke-dasharray:",
        Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.4),
        size=14, bold=True, colore=AZZURRO)

    steps2 = [
        "Misura la\nlunghezza totale\ndella path SVG",
        "Imposta\ndasharray = L\ndashoffset = L",
        "Linea\ninvisibile\n(offset = gap)",
        "Transizione D3\noffset -> 0\ndurata 1500ms",
        "Linea si disegna\nda sx a dx\nnel tempo",
    ]
    for i, label in enumerate(steps2):
        col = RGBColor(0x0A, 0x3D, 0x2B) if i == 4 else ACCENT
        x_n = Inches(0.4) + i * Inches(2.5)
        nodo(sl, label, x_n, Inches(3.72), Inches(2.2), Inches(1.1), col=col, size=11)
        if i < 4:
            freccia_txt(sl, x_n + Inches(2.25), Inches(4.05))

    txt(sl, "Effetto: la linea si \"disegna\" progressivamente — guida l'occhio sull'asse temporale.",
        Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.45),
        size=13, italic=True, colore=GIALLO)

    rettangolo(sl, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.75), ACCENT)
    txt(sl, "Parallel coordinates — scelta analoga: 300 utenti -> 200 campionati (shuffle casuale) "
            "per evitare overplotting e blocco del browser.",
        Inches(0.65), Inches(5.67), Inches(12.0), Inches(0.62),
        size=13, colore=GRIGIO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Flusso filtri (coordinated views)
# ════════════════════════════════════════════════════════════════════════════════

def slide_coordinated(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — Flusso dei Filtri")
    sep_orizz(sl)

    txt(sl, "Ogni filtro si propaga simultaneamente a tutte e 7 le visualizzazioni:",
        Inches(0.5), Inches(1.28), Inches(12.3), Inches(0.4),
        size=13, italic=True, colore=GRIGIO)

    flusso = [
        "Utente\ninteragisce\n(click / slider)",
        "interactions.js\nrileva l'evento",
        "filterData()\ncrea il subset",
        "updateAllCharts()\nchiama update()",
        "7 grafici\nsi ridisegnano",
    ]
    for i, label in enumerate(flusso):
        col = RGBColor(0x0A, 0x3D, 0x2B) if i == 4 else ACCENT
        x_n = Inches(0.4) + i * Inches(2.5)
        nodo(sl, label, x_n, Inches(1.82), Inches(2.2), Inches(1.05), col=col, size=12)
        if i < 4:
            freccia_txt(sl, x_n + Inches(2.25), Inches(2.1))

    # filtri disponibili (colonna sx)
    txt(sl, "I filtri disponibili:", Inches(0.5), Inches(3.15),
        Inches(6.0), Inches(0.4), size=14, bold=True, colore=AZZURRO)

    filtri = [
        ("Cluster",      "multi-select C0 ... C4"),
        ("Tipo utente",  "tutti  /  solo insider  /  solo normali"),
        ("Score range",  "doppio slider — input: aggiorna label  |  change: applica filtro"),
        ("Reset",        "ripristina il dataset completo"),
    ]
    for i, (nome, desc) in enumerate(filtri):
        y_f = Inches(3.65) + i * Inches(0.78)
        rettangolo(sl, Inches(0.5), y_f, Inches(6.0), Inches(0.72), ACCENT)
        rettangolo(sl, Inches(0.5), y_f, Inches(0.07), Inches(0.72), AZZURRO)
        txt(sl, nome, Inches(0.7), y_f + Inches(0.08),
            Inches(1.6), Inches(0.55), size=13, bold=True, colore=AZZURRO)
        txt(sl, desc, Inches(2.4), y_f + Inches(0.08),
            Inches(4.0), Inches(0.55), size=12, colore=GRIGIO)

    # architettura hub-and-spoke (colonna dx)
    txt(sl, "Architettura hub-and-spoke:", Inches(7.1), Inches(3.15),
        Inches(5.8), Inches(0.4), size=14, bold=True, colore=AZZURRO)

    rettangolo(sl, Inches(9.1), Inches(3.75), Inches(2.0), Inches(0.7),
               RGBColor(0x0F, 0x3A, 0x60))
    rettangolo(sl, Inches(9.1), Inches(3.75), Inches(0.1), Inches(0.7), ROSSO)
    txt(sl, "interactions.js  (hub)", Inches(9.25), Inches(3.83),
        Inches(1.75), Inches(0.55), size=12, bold=True, colore=ROSSO)

    spokes = [
        (Inches(7.1),  Inches(3.75), "univariate"),
        (Inches(7.1),  Inches(4.55), "bivariate"),
        (Inches(7.1),  Inches(5.35), "trivariate"),
        (Inches(11.2), Inches(3.75), "temporal"),
        (Inches(11.2), Inches(4.55), "parallel-coord"),
        (Inches(11.2), Inches(5.35), "multivariate"),
    ]
    for x_s, y_s, nome in spokes:
        rettangolo(sl, x_s, y_s, Inches(1.8), Inches(0.6), SFONDO)
        txt(sl, nome, x_s + Inches(0.1), y_s + Inches(0.1),
            Inches(1.6), Inches(0.4), size=11, colore=GRIGIO)

    txt(sl, "I moduli non si conoscono tra loro — ognuno espone solo init() e update().",
        Inches(7.1), Inches(6.2), Inches(5.8), Inches(0.65),
        size=12, italic=True, colore=GIALLO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — I 5 cluster
# ════════════════════════════════════════════════════════════════════════════════

def slide_cluster(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Progetto Finale — I 5 Cluster")
    sep_orizz(sl)

    txt(sl, "K-Means k=5 sui pattern comportamentali — i numeri 0-4 sono arbitrari",
        Inches(0.5), Inches(1.28), Inches(12.3), Inches(0.4),
        size=13, italic=True, colore=GRIGIO)

    colori_cl = [GRIGIO, AZZURRO, VERDE, AZZURRO, ROSSO]
    clusters = [
        ("C0", "Poco attivi",     "0",  "pochi logon, bassa attivita' generale"),
        ("C1", "After-hour",      "0",  "after_hour_ratio il piu' alto — non sono insider"),
        ("C2", "Alta intensita'", "0",  "score medio 1.46 — il piu' alto del dataset, zero insider"),
        ("C3", "File-heavy",      "0",  "molte operazioni su file, attivi in orario"),
        ("C4", "Insider",         "15", "TUTTI i 15 insider reali — 25% del cluster"),
    ]
    col_w_cl = [Inches(1.2), Inches(2.4), Inches(1.2), Inches(7.6)]
    x0, riga_h, y0 = Inches(0.5), Inches(0.65), Inches(1.82)

    for j, (h_txt, cw) in enumerate(zip(["Cluster","Profilo","Insider","Note"], col_w_cl)):
        x_c = x0 + sum(col_w_cl[:j])
        cella = rettangolo(sl, x_c, y0, cw, riga_h, ACCENT)
        tf = cella.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h_txt
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = AZZURRO

    for i, (cn, profilo, ins, note) in enumerate(clusters):
        bg_r = RGBColor(0x10, 0x10, 0x25) if i % 2 == 0 else SFONDO
        for j, (val, cw) in enumerate(zip([cn, profilo, ins, note], col_w_cl)):
            x_c = x0 + sum(col_w_cl[:j])
            cella = rettangolo(sl, x_c, y0 + riga_h * (i + 1), cw, riga_h, bg_r)
            tf = cella.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j < 3 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(13)
            run.font.bold = (i == 4)
            run.font.color.rgb = (colori_cl[i] if j == 0
                                  else ROSSO if (i == 4 and j == 2)
                                  else GRIGIO)

    rettangolo(sl, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.85), ACCENT)
    txt(sl, "Insight: score medio C4 = 0.92  <  C2 = 1.46.  "
            "Alto score NON identifica gli insider — conta il pattern comportamentale.",
        Inches(0.65), Inches(6.27), Inches(12.0), Inches(0.7),
        size=13, bold=True, colore=ROSSO)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusioni
# ════════════════════════════════════════════════════════════════════════════════

def slide_conclusioni(prs):
    sl = nuova_slide(prs)
    sfondo_slide(sl)
    barra_titolo(sl, "Conclusioni")
    sep_orizz(sl)

    txt(sl, "Cosa ho imparato su D3.js", Inches(0.5), Inches(1.3),
        Inches(5.9), Inches(0.45), size=16, bold=True, colore=AZZURRO)
    bullets(sl, [
        "Il data join (select / data / enter / exit) e' il cuore di D3",
        "Le scale sono funzioni pure: dato  ->  pixel",
        "Ogni transizione: .duration().ease().attr()",
        "d3.rollup / d3.bin trasformano dati in strutture visualizzabili",
        "I grafici si coordinano tramite architettura, non magia",
        "La margin convention separa logica e layout SVG",
    ], Inches(0.5), Inches(1.82), Inches(5.9), Inches(3.4), size=14)

    txt(sl, "Intermedio vs Finale", Inches(7.1), Inches(1.3),
        Inches(5.8), Inches(0.45), size=16, bold=True, colore=AZZURRO)
    tabella(sl,
        ["", "Intermedio", "Finale"],
        [
            ("Dati",          "JSON, 10 record",     "CSV, 300 record"),
            ("Caricamento",   "d3.json()",           "d3.csv() x2 + join"),
            ("Scale",         "scaleLinear x2",      "bin / sqrt / point / band"),
            ("Interazione",   "click su glifo",      "filtri coordinated"),
            ("Moduli",        "1 file app.js",        "10 moduli separati"),
            ("Aggiornamento", "transizione singola",  "update() su tutti"),
        ],
        Inches(7.1), Inches(1.82),
        [Inches(1.8), Inches(1.9), Inches(2.05)], Inches(0.6))

    sep_orizz(sl, Inches(5.7))
    txt(sl, "D3 non disegna grafici — fornisce gli strumenti per trasformare "
            "dati in elementi visivi controllabili.",
        Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.9),
        size=16, italic=True, colore=BIANCO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════════════════════════

def genera():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_copertina(prs)           # 1
    slide_sep_intermedio(prs)      # 2
    slide_intermedio_overview(prs) # 3
    slide_intermedio_d3(prs)       # 4
    slide_intermedio_glifo(prs)    # 5
    slide_sep_finale(prs)          # 6
    slide_dataset(prs)             # 7
    slide_dati_csv(prs)            # 8  — struttura dei due CSV
    slide_architettura(prs)        # 9  — architettura moduli JS
    slide_architettura_dati(prs)   # 10 — colonna -> grafico
    slide_graf_istogramma(prs)     # 11 — grafico 1
    slide_graf_scatter(prs)        # 12 — grafico 2
    slide_graf_bubble(prs)         # 13 — grafico 3
    slide_graf_line(prs)           # 14 — grafico 4
    slide_graf_parallel(prs)       # 15 — grafico 5
    slide_graf_radar(prs)          # 16 — grafico 6
    slide_graf_heatmap(prs)        # 17 — grafico 7
    slide_flusso_viz(prs)          # 18 — scale D3 riepilogo
    slide_flusso_temporale(prs)    # 19 — animazione linee
    slide_coordinated(prs)         # 20 — filtri coordinated
    slide_cluster(prs)             # 21 — i 5 cluster
    slide_conclusioni(prs)         # 22

    out = "_docs/presentazione_infovis.pptx"
    prs.save(out)
    print(f"Salvato: {out}  ({len(prs.slides)} slide)")


if __name__ == "__main__":
    genera()
